"""
The Arena Protocol — Investor Pitch Deck Generator

Generates a 12-slide professional PowerPoint pitch deck.
Dark theme, navy + gold color scheme, no emoji, no clip art.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════
# DESIGN SYSTEM
# ═══════════════════════════════════════════════════

# Colors
NAVY_DARK   = RGBColor(0x08, 0x0E, 0x1A)   # Background
NAVY_MID    = RGBColor(0x12, 0x1D, 0x33)   # Card bg
NAVY_LIGHT  = RGBColor(0x1B, 0x2A, 0x4A)   # Borders, secondary
GOLD        = RGBColor(0xD4, 0xA5, 0x37)   # Primary accent
GOLD_LIGHT  = RGBColor(0xE8, 0xC5, 0x6D)   # Secondary accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xB0, 0xB8, 0xC8)
MED_GRAY    = RGBColor(0x78, 0x82, 0x96)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
RED         = RGBColor(0xEF, 0x44, 0x44)
BLUE        = RGBColor(0x3B, 0x82, 0xF6)

# Fonts
FONT_HEADING = 'Calibri'
FONT_BODY    = 'Calibri'

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════

def set_slide_bg(slide, color=NAVY_DARK):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY, line_spacing=1.2):
    """Add a text box with single-run formatting."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=14,
                      color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                      line_spacing=1.5, font_name=FONT_BODY):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Handle tuples for per-line formatting
        if isinstance(line, tuple):
            txt, fmt = line
            p.text = txt
            p.font.size = Pt(fmt.get('size', font_size))
            p.font.color.rgb = fmt.get('color', color)
            p.font.bold = fmt.get('bold', bold)
            p.font.name = fmt.get('font', font_name)
            p.alignment = fmt.get('align', alignment)
        else:
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.name = font_name
            p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_rect(slide, left, top, width, height, fill_color=NAVY_MID,
             border_color=None, border_width=Pt(1)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_circle(slide, left, top, size, fill_color=GOLD):
    """Add a circle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=GOLD, width=Pt(2)):
    """Add a line connector."""
    connector = slide.shapes.add_connector(
        1,  # straight
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_slide_number(slide, num, total=12):
    """Add slide number in bottom right."""
    add_text_box(slide, 11.8, 7.0, 1.2, 0.4, f"{num} / {total}",
                 font_size=9, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


def add_section_title(slide, title, subtitle=None):
    """Standard section title + subtitle layout."""
    # Gold accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(0.7), Inches(0.06), Inches(0.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    add_text_box(slide, 1.1, 0.6, 10, 0.7, title,
                 font_size=32, color=WHITE, bold=True,
                 font_name=FONT_HEADING)
    if subtitle:
        add_text_box(slide, 1.1, 1.2, 10, 0.5, subtitle,
                     font_size=16, color=LIGHT_GRAY)


def add_stat_card(slide, left, top, width, height, value, label,
                  value_color=GOLD, bg_color=NAVY_MID):
    """Add a stat card with large value and label."""
    add_rect(slide, left, top, width, height, fill_color=bg_color,
             border_color=NAVY_LIGHT)
    add_text_box(slide, left + 0.2, top + 0.25, width - 0.4, 0.6,
                 value, font_size=28, color=value_color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + 0.2, top + 0.85, width - 0.4, 0.4,
                 label, font_size=11, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════

def slide_01_title(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Gold horizontal rule
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(2.8), Inches(6.3), Pt(3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    add_text_box(slide, 3.5, 2.0, 6.3, 0.8, 'THE ARENA',
                 font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

    add_text_box(slide, 3.0, 3.1, 7.3, 0.6,
                 'The Adversarial Execution Protocol for AI Agents',
                 font_size=22, color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, 3.5, 4.0, 6.3, 0.8,
                 'Sealed-bid auctions. Staked execution. Independent verification.\nOn-chain slashing. Zero trust required.',
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.6)

    # Bottom tag
    add_text_box(slide, 3.5, 6.0, 6.3, 0.4,
                 'Pre-Seed  |  $250K Raise  |  Base (Ethereum L2)',
                 font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 1)


def slide_02_problem(prs):
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, 'The Problem')

    add_text_box(slide, 1.1, 1.6, 11, 0.6,
                 'AI agents manage over $2B in DeFi capital today with zero accountability infrastructure.',
                 font_size=18, color=LIGHT_GRAY)

    # 3 problem cards
    cards = [
        ('No Accountability',
         'AI agents execute trades, manage vaults, and move millions with no performance bonds, no staking, and no consequences for failure.',
         '$2B+', 'in DeFi managed by AI agents'),
        ('No Verification',
         'Agent outputs are accepted on faith. There is no independent verification layer, no schema validation, and no audit trail.',
         '0%', 'of agent outputs are independently verified'),
        ('No Skin in the Game',
         'Agents earn fees regardless of outcome quality. Bad actors profit the same as competent ones. Misaligned incentives at every level.',
         '$0', 'at risk when agents deliver poor results'),
    ]

    for i, (title, desc, stat, stat_label) in enumerate(cards):
        x = 1.1 + i * 3.8
        add_rect(slide, x, 2.5, 3.5, 4.2, fill_color=NAVY_MID,
                 border_color=NAVY_LIGHT)
        # Red top accent
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(2.5), Inches(3.5), Pt(4)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RED
        bar.line.fill.background()

        add_text_box(slide, x + 0.3, 2.8, 2.9, 0.4, title,
                     font_size=18, color=WHITE, bold=True)
        add_text_box(slide, x + 0.3, 3.3, 2.9, 1.8, desc,
                     font_size=12, color=LIGHT_GRAY, line_spacing=1.5)

        # Stat at bottom
        add_text_box(slide, x + 0.3, 5.3, 2.9, 0.5, stat,
                     font_size=28, color=RED, bold=True)
        add_text_box(slide, x + 0.3, 5.9, 2.9, 0.4, stat_label,
                     font_size=10, color=MED_GRAY)

    add_slide_number(slide, 2)


def slide_03_solution(prs):
    """Slide 3: The Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, 'The Solution')

    add_text_box(slide, 1.1, 1.6, 11, 0.6,
                 'An adversarial execution protocol that forces accountability through economic incentives.',
                 font_size=18, color=LIGHT_GRAY)

    # 4 solution cards
    solutions = [
        ('Sealed-Bid Auctions',
         'Agents compete on price, stake, and ETA via commit-reveal sealed bids. Scoring: (stake x reputation) / price. No collusion.',
         GOLD),
        ('Staked Execution',
         'Winning agents lock capital as a performance bond. Minimum 10% of bounty. Skin in the game for every task.',
         GREEN),
        ('Independent Verification',
         'VRF-selected verifiers independently validate outputs against acceptance criteria. No rubber-stamping.',
         BLUE),
        ('On-Chain Slashing',
         'Failed delivery triggers automatic slashing at 5 severity tiers: 15% to 100%. 30-day post-completion slash window.',
         RED),
    ]

    for i, (title, desc, accent) in enumerate(solutions):
        x = 0.8 + i * 3.1
        add_rect(slide, x, 2.5, 2.85, 4.2, fill_color=NAVY_MID,
                 border_color=NAVY_LIGHT)
        # Accent top bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(2.5), Inches(2.85), Pt(4)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        add_text_box(slide, x + 0.25, 2.8, 2.35, 0.4, title,
                     font_size=16, color=WHITE, bold=True)
        add_text_box(slide, x + 0.25, 3.3, 2.35, 2.8, desc,
                     font_size=12, color=LIGHT_GRAY, line_spacing=1.5)

    add_slide_number(slide, 3)


def slide_04_how_it_works(prs):
    """Slide 4: How It Works"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, 'How It Works')

    steps = [
        ('1', 'POST', 'Task poster defines acceptance criteria, sets bounty, deposits USDC into escrow.'),
        ('2', 'BID', 'Agents submit sealed bids (stake + price + ETA). Commit-reveal prevents front-running.'),
        ('3', 'VERIFY', 'VRF-selected verifiers independently validate output against criteria. Consensus required.'),
        ('4', 'SETTLE', 'Approved: agent paid, stake returned. Failed: stake slashed, bond held for 30 days.'),
    ]

    y_top = 2.3
    for i, (num, label, desc) in enumerate(steps):
        x = 0.8 + i * 3.15

        # Step circle
        add_circle(slide, x + 0.85, y_top, 0.7, fill_color=GOLD)
        add_text_box(slide, x + 0.85, y_top + 0.1, 0.7, 0.6, num,
                     font_size=24, color=NAVY_DARK, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Arrow connector (except last)
        if i < 3:
            add_line(slide, x + 1.7, y_top + 0.35, x + 3.0, y_top + 0.35,
                     color=GOLD, width=Pt(2))

        # Label
        add_text_box(slide, x, y_top + 0.9, 2.7, 0.4, label,
                     font_size=20, color=GOLD, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Description card
        add_rect(slide, x, y_top + 1.4, 2.7, 2.5, fill_color=NAVY_MID,
                 border_color=NAVY_LIGHT)
        add_text_box(slide, x + 0.2, y_top + 1.6, 2.3, 2.2, desc,
                     font_size=13, color=LIGHT_GRAY, line_spacing=1.5)

    # Bottom callout
    add_rect(slide, 0.8, 6.5, 11.7, 0.7, fill_color=NAVY_MID,
             border_color=GOLD)
    add_text_box(slide, 1.0, 6.55, 11.3, 0.6,
                 'Scoring Formula:  score = (stake x (reputation + 1) x 1e18) / price    |    All on-chain. All verifiable. All auditable.',
                 font_size=13, color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 4)


def slide_05_market(prs):
    """Slide 5: Market Size"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, 'Market Opportunity')

    # TAM / SAM / SOM boxes
    markets = [
        ('TAM', '$47B', 'AI Agent Market (2027)',
         'Autonomous AI agents managing capital, executing trades, and making decisions across DeFi, TradFi, and enterprise.'),
        ('SAM', '$8.2B', 'DeFi AI Agent Activity',
         'AI agents actively participating in DeFi protocols: vault management, MEV, liquidations, trading, and risk assessment.'),
        ('SOM', '$180M', 'Protocol Fee Capture (Y3)',
         'Arena protocol fees on staked execution: settlement, slashing, disputes, insurance, and SaaS data intelligence.'),
    ]

    for i, (tag, value, label, desc) in enumerate(markets):
        x = 1.1 + i * 3.8
        add_rect(slide, x, 2.0, 3.5, 3.0, fill_color=NAVY_MID,
                 border_color=NAVY_LIGHT)
        add_text_box(slide, x + 0.3, 2.15, 1.2, 0.35, tag,
                     font_size=14, color=GOLD, bold=True)
        add_text_box(slide, x + 0.3, 2.5, 2.9, 0.6, value,
                     font_size=36, color=WHITE, bold=True)
        add_text_box(slide, x + 0.3, 3.1, 2.9, 0.35, label,
                     font_size=13, color=GOLD_LIGHT)
        add_text_box(slide, x + 0.3, 3.5, 2.9, 1.2, desc,
                     font_size=11, color=LIGHT_GRAY, line_spacing=1.4)

    # DeFi TVL context
    add_rect(slide, 1.1, 5.3, 11.2, 1.8, fill_color=NAVY_MID,
             border_color=NAVY_LIGHT)

    add_text_box(slide, 1.4, 5.5, 10, 0.4, 'MARKET TAILWINDS',
                 font_size=14, color=GOLD, bold=True)

    tailwinds = [
        'DeFi TVL exceeded $250B in 2025 with increasing AI-managed capital across lending, DEXs, and yield protocols.',
        'Enterprise adoption of AI agents for financial operations is growing 40%+ annually.',
        'No existing protocol provides adversarial execution infrastructure. The Arena creates a new category.',
    ]
    for j, tw in enumerate(tailwinds):
        add_text_box(slide, 1.7, 5.95 + j * 0.4, 10, 0.35,
                     f'   {tw}', font_size=11, color=LIGHT_GRAY)

    add_slide_number(slide, 5)


def slide_06_business_model(prs):
    """Slide 6: Business Model — 5 Revenue Streams"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, 'Business Model', '5 protocol-native revenue streams. No token dependency.')

    streams = [
        ('2.5%', 'Settlement Fee', 'On every successfully completed task. Applied to total bounty at settlement.', GOLD),
        ('10%', 'Slash Fee', 'Protocol cut of slashed agent stakes. 5 severity tiers from 15% to 100%.', RED),
        ('5%', 'Dispute Fee', 'Charged on every dispute filed. Funds arbitrator compensation and protocol treasury.', BLUE),
        ('1%', 'Insurance Cut', 'Protocol fee on insurance premiums. Agents can buy coverage against slashing risk.', GREEN),
        ('SaaS', 'Agent Reliability Index', 'Subscription data product: agent scoring, task outcome predictions, risk analytics.', GOLD_LIGHT),
    ]

    for i, (rate, name, desc, accent) in enumerate(streams):
        x = 0.6 + i * 2.5
        add_rect(slide, x, 2.2, 2.25, 4.5, fill_color=NAVY_MID,
                 border_color=NAVY_LIGHT)
        # Accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(2.2), Inches(2.25), Pt(4)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        add_text_box(slide, x + 0.15, 2.5, 1.95, 0.6, rate,
                     font_size=32, color=accent, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.15, 3.15, 1.95, 0.4, name,
                     font_size=14, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.15, 3.6, 1.95, 2.5, desc,
                     font_size=11, color=LIGHT_GRAY, line_spacing=1.5,
                     alignment=PP_ALIGN.CENTER)

    # Bottom note
    add_text_box(slide, 1.1, 7.0, 11, 0.4,
                 'Blended take rate: 3.7% at scale  |  Fee phasing: M1-3 free, M4-6 at 1%, M7+ full 2.5%',
                 font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 6)


def slide_07_traction(prs):
    """Slide 7: Traction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, 'Traction', 'Built, tested, and audit-ready.')

    # Stat cards row 1
    stats = [
        ('819', 'Tests Passing', GOLD),
        ('10', 'Solidity Contracts', GOLD),
        ('5', 'Autonomous Agents', GOLD),
        ('100%', 'Slither Clean', GREEN),
    ]
    for i, (val, label, color) in enumerate(stats):
        x = 0.8 + i * 3.1
        add_stat_card(slide, x, 2.0, 2.7, 1.4, val, label,
                      value_color=color)

    # Built items
    items = [
        ('Protocol Contracts', '9 satellite contracts + ArenaCore. Sealed-bid auctions, VRF verification, 5-tier slashing, continuous contracts, syndicates, delegation, insurance, arbitration, OFAC compliance.'),
        ('SDK + Agent Framework', 'TypeScript SDK with typed ABIs for all 10 contracts. Audit Agent, Verifier Agent, Risk Agent, Task Poster Bot, Agent Orchestrator.'),
        ('Frontend + Landing', 'Next.js dashboard (wagmi v2 + RainbowKit) with task lifecycle, agent leaderboard, create-task flow. Marketing landing page.'),
        ('Deployment Pipeline', 'Hardhat deployment script: 11 contracts, satellite linking, USDC whitelisting, Timelock ownership transfer. Tested on local Hardhat network.'),
    ]

    for i, (title, desc) in enumerate(items):
        y = 3.7 + i * 0.9
        add_rect(slide, 0.8, y, 11.7, 0.8, fill_color=NAVY_MID,
                 border_color=NAVY_LIGHT)
        add_text_box(slide, 1.1, y + 0.05, 2.5, 0.35, title,
                     font_size=13, color=GOLD, bold=True)
        add_text_box(slide, 3.6, y + 0.05, 8.5, 0.7, desc,
                     font_size=11, color=LIGHT_GRAY, line_spacing=1.3)

    add_slide_number(slide, 7)


def slide_08_financials(prs):
    """Slide 8: Financial Projections"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, 'Financial Projections',
                      'Base case from 24-month P&L model. All formulas reference assumption inputs.')

    # M12 column
    add_rect(slide, 0.8, 2.2, 5.8, 5.0, fill_color=NAVY_MID,
             border_color=NAVY_LIGHT)
    add_text_box(slide, 1.1, 2.35, 5.2, 0.4, 'MONTH 12',
                 font_size=18, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    m12_stats = [
        ('Monthly Revenue', '$55,788'),
        ('Monthly GMV', '$984,693'),
        ('Monthly Tasks', '371'),
        ('Cumulative Tasks', '1,569'),
        ('Cumulative GMV', '$3.8M'),
        ('SaaS ARR', '$304,570'),
        ('Total ARR', '$669,450'),
        ('Net Margin', '96.2%'),
    ]
    for i, (label, value) in enumerate(m12_stats):
        y = 2.85 + i * 0.52
        add_text_box(slide, 1.3, y, 3.0, 0.35, label,
                     font_size=12, color=LIGHT_GRAY)
        add_text_box(slide, 4.3, y, 2.0, 0.35, value,
                     font_size=13, color=WHITE, bold=True,
                     alignment=PP_ALIGN.RIGHT)

    # M24 column
    add_rect(slide, 6.9, 2.2, 5.8, 5.0, fill_color=NAVY_MID,
             border_color=GOLD)
    add_text_box(slide, 7.2, 2.35, 5.2, 0.4, 'MONTH 24',
                 font_size=18, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    m24_stats = [
        ('Monthly Revenue', '$1,424,650'),
        ('Monthly GMV', '$38.6M'),
        ('Monthly Tasks', '8,369'),
        ('Cumulative Tasks', '36,398'),
        ('Cumulative GMV', '$145.2M'),
        ('SaaS ARR', '$2.7M'),
        ('Total ARR', '$17.1M'),
        ('Net Margin', '99.2%'),
    ]
    for i, (label, value) in enumerate(m24_stats):
        y = 2.85 + i * 0.52
        add_text_box(slide, 7.4, y, 3.0, 0.35, label,
                     font_size=12, color=LIGHT_GRAY)
        add_text_box(slide, 10.4, y, 2.0, 0.35, value,
                     font_size=13, color=WHITE, bold=True,
                     alignment=PP_ALIGN.RIGHT)

    add_slide_number(slide, 8)


def slide_09_saas(prs):
    """Slide 9: SaaS Opportunity"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, 'SaaS Opportunity',
                      'Agent Reliability Index — proprietary data product built on on-chain execution history.')

    # Product description
    add_rect(slide, 0.8, 2.0, 7.0, 2.2, fill_color=NAVY_MID,
             border_color=NAVY_LIGHT)
    add_text_box(slide, 1.1, 2.15, 6.4, 0.35, 'AGENT RELIABILITY INDEX',
                 font_size=16, color=GOLD, bold=True)

    features = [
        'AI agent performance scoring based on on-chain execution history',
        'Task outcome prediction models trained on protocol data',
        'Risk analytics: slash probability, failure correlations, anomaly detection',
        'Enterprise API: query agent reliability before delegation or hiring',
    ]
    for j, feat in enumerate(features):
        add_text_box(slide, 1.3, 2.55 + j * 0.35, 6.2, 0.3,
                     f'   {feat}', font_size=11, color=LIGHT_GRAY)

    # Pricing tiers
    tiers = [
        ('Basic', '$500/mo', '60%', 'API access, basic scoring, 1,000 queries/mo'),
        ('Pro', '$2,000/mo', '30%', 'Full analytics, predictions, 10,000 queries/mo'),
        ('Enterprise', '$8,000/mo', '10%', 'Custom models, SLAs, unlimited queries, dedicated support'),
    ]

    add_text_box(slide, 8.1, 2.15, 4.5, 0.35, 'PRICING TIERS',
                 font_size=14, color=GOLD, bold=True)

    for j, (tier, price, mix, desc) in enumerate(tiers):
        y = 2.6 + j * 0.55
        add_rect(slide, 8.1, y, 4.5, 0.5, fill_color=NAVY_MID,
                 border_color=NAVY_LIGHT)
        add_text_box(slide, 8.3, y + 0.05, 1.2, 0.3, tier,
                     font_size=12, color=WHITE, bold=True)
        add_text_box(slide, 9.5, y + 0.05, 1.0, 0.3, price,
                     font_size=12, color=GOLD)
        add_text_box(slide, 10.5, y + 0.05, 0.5, 0.3, mix,
                     font_size=11, color=MED_GRAY)
        add_text_box(slide, 11.0, y + 0.05, 1.5, 0.3, '',
                     font_size=10, color=LIGHT_GRAY)

    add_text_box(slide, 8.1, 4.3, 4.5, 0.3,
                 'Blended ARPU: $1,700/mo',
                 font_size=12, color=GOLD_LIGHT, bold=True)

    # Comparables
    add_rect(slide, 0.8, 4.6, 11.7, 2.5, fill_color=NAVY_MID,
             border_color=NAVY_LIGHT)
    add_text_box(slide, 1.1, 4.75, 10, 0.35, 'COMPARABLE COMPANIES',
                 font_size=14, color=GOLD, bold=True)

    comps = [
        ('Chainalysis', 'Blockchain Analytics', '$200M+ ARR'),
        ('Nansen', 'On-Chain Analytics', '$30M+ ARR'),
        ('Dune Analytics', 'Data Dashboards', '$20M+ ARR'),
        ('Gauntlet', 'DeFi Risk Modeling', '$15M+ ARR'),
        ('The Arena (Y2)', 'AI Agent Intelligence', '$2.7M ARR'),
    ]

    for j, (company, category, arr) in enumerate(comps):
        y = 5.2 + j * 0.35
        is_arena = j == 4
        add_text_box(slide, 1.3, y, 2.5, 0.3, company,
                     font_size=12, color=GOLD if is_arena else WHITE,
                     bold=is_arena)
        add_text_box(slide, 4.0, y, 3.5, 0.3, category,
                     font_size=12, color=GOLD_LIGHT if is_arena else LIGHT_GRAY)
        add_text_box(slide, 8.0, y, 3.0, 0.3, arr,
                     font_size=12, color=GOLD if is_arena else WHITE,
                     bold=is_arena, alignment=PP_ALIGN.RIGHT)

    add_slide_number(slide, 9)


def slide_10_gtm(prs):
    """Slide 10: Go-to-Market"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, 'Go-to-Market Timeline')

    phases = [
        ('M1-2', 'Testnet + Audit', [
            'Smart contract audit with formal verification',
            'Base Sepolia testnet deployment live',
            'Seed tasks funded, initial agent onboarding',
            'Bug bounty program launched',
        ], GOLD),
        ('M3-4', 'Mainnet Launch', [
            'Base mainnet deployment + Timelock governance',
            '50+ tasks posted, first agent earnings',
            'Fee phasing begins at 1% settlement rate',
            'Insurance marketplace opens',
        ], GREEN),
        ('M5-8', 'Scale + SaaS', [
            'Full 2.5% fee schedule activates at M7',
            'Agent Reliability Index SaaS product launches',
            'Continuous contract support live',
            'Community growth + partnership outreach',
        ], BLUE),
        ('M9-12', 'Growth + Fundraise', [
            '1,000+ cumulative tasks milestone',
            'SaaS ARR exceeds $300K',
            'Multi-chain expansion planning',
            'Series A preparation and outreach',
        ], GOLD_LIGHT),
    ]

    for i, (period, title, items, accent) in enumerate(phases):
        x = 0.6 + i * 3.15
        # Phase card
        add_rect(slide, x, 2.0, 2.9, 5.0, fill_color=NAVY_MID,
                 border_color=NAVY_LIGHT)
        # Accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(2.0), Inches(2.9), Pt(4)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        add_text_box(slide, x + 0.2, 2.25, 2.5, 0.35, period,
                     font_size=20, color=accent, bold=True)
        add_text_box(slide, x + 0.2, 2.65, 2.5, 0.35, title,
                     font_size=14, color=WHITE, bold=True)

        # Bullet items
        for j, item in enumerate(items):
            add_text_box(slide, x + 0.2, 3.15 + j * 0.6, 2.5, 0.55,
                         item, font_size=11, color=LIGHT_GRAY,
                         line_spacing=1.3)

        # Connector arrow (except last)
        if i < 3:
            arrow_x = x + 3.0
            add_line(slide, arrow_x, 4.5, arrow_x + 0.2, 4.5,
                     color=MED_GRAY, width=Pt(1.5))

    add_slide_number(slide, 10)


def slide_11_use_of_funds(prs):
    """Slide 11: Use of Funds"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, 'Use of Funds', '$250K pre-seed raise')

    allocations = [
        ('Team (6 months)', '$60K', '24%', 'Core team compensation runway'),
        ('Marketing & Growth', '$40K', '16%', 'Community building, content, partnerships'),
        ('SaaS Development', '$30K', '12%', 'Agent Reliability Index data product build'),
        ('Reserve', '$30K', '12%', 'Operating buffer and contingency'),
        ('Seed Task Bounties', '$20K', '8%', 'Bootstrap initial protocol activity'),
        ('Bug Bounty Program', '$20K', '8%', 'Security researcher rewards'),
        ('Smart Contract Audit', '$15K', '6%', 'Formal verification + third-party audit'),
        ('Agent Incentives', '$15K', '6%', 'Early agent onboarding rewards'),
        ('Legal & Compliance', '$10K', '4%', 'Entity formation, regulatory review'),
        ('Infrastructure', '$10K', '4%', 'RPC nodes, IPFS, monitoring, hosting'),
    ]

    # Header row
    add_rect(slide, 0.8, 2.2, 11.7, 0.5, fill_color=NAVY_LIGHT)
    add_text_box(slide, 1.1, 2.25, 3.5, 0.35, 'CATEGORY',
                 font_size=11, color=GOLD, bold=True)
    add_text_box(slide, 4.8, 2.25, 1.2, 0.35, 'AMOUNT',
                 font_size=11, color=GOLD, bold=True,
                 alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, 6.2, 2.25, 0.8, 0.35, '%',
                 font_size=11, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 7.2, 2.25, 5.0, 0.35, 'PURPOSE',
                 font_size=11, color=GOLD, bold=True)

    for i, (cat, amount, pct, purpose) in enumerate(allocations):
        y = 2.8 + i * 0.43
        bg = NAVY_MID if i % 2 == 0 else NAVY_DARK
        add_rect(slide, 0.8, y, 11.7, 0.4, fill_color=bg)

        # Bar chart visual
        pct_val = float(pct.replace('%', '')) / 100
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(y), Inches(11.7 * pct_val), Pt(3)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = GOLD
        bar.line.fill.background()

        add_text_box(slide, 1.1, y + 0.02, 3.5, 0.3, cat,
                     font_size=12, color=WHITE)
        add_text_box(slide, 4.8, y + 0.02, 1.2, 0.3, amount,
                     font_size=12, color=WHITE, bold=True,
                     alignment=PP_ALIGN.RIGHT)
        add_text_box(slide, 6.2, y + 0.02, 0.8, 0.3, pct,
                     font_size=11, color=GOLD,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, 7.2, y + 0.02, 5.0, 0.3, purpose,
                     font_size=11, color=LIGHT_GRAY)

    # Total row
    y_total = 2.8 + len(allocations) * 0.43 + 0.1
    add_rect(slide, 0.8, y_total, 11.7, 0.45, fill_color=NAVY_LIGHT,
             border_color=GOLD)
    add_text_box(slide, 1.1, y_total + 0.05, 3.5, 0.3, 'TOTAL',
                 font_size=13, color=GOLD, bold=True)
    add_text_box(slide, 4.8, y_total + 0.05, 1.2, 0.3, '$250K',
                 font_size=13, color=GOLD, bold=True,
                 alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, 6.2, y_total + 0.05, 0.8, 0.3, '100%',
                 font_size=13, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 11)


def slide_12_team(prs):
    """Slide 12: Team / Contact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Title centered
    add_text_box(slide, 3.0, 1.5, 7.3, 0.8, 'THE ARENA',
                 font_size=40, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

    # Gold rule
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.5), Inches(2.3), Inches(4.3), Pt(3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    add_text_box(slide, 3.0, 2.5, 7.3, 0.5,
                 'Adversarial Execution Protocol for AI Agents',
                 font_size=18, color=GOLD_LIGHT, alignment=PP_ALIGN.CENTER)

    # Key metrics summary
    summary_items = [
        ('819 tests', '10 contracts', 'Slither clean'),
        ('$250K raise', '5 autonomous agents', '12-month path to $669K ARR'),
    ]
    for row_i, row in enumerate(summary_items):
        for col_i, item in enumerate(row):
            x = 2.8 + col_i * 2.8
            y = 3.5 + row_i * 0.5
            add_text_box(slide, x, y, 2.5, 0.4, item,
                         font_size=13, color=LIGHT_GRAY,
                         alignment=PP_ALIGN.CENTER)

    # Contact section
    add_rect(slide, 3.5, 4.8, 6.3, 1.8, fill_color=NAVY_MID,
             border_color=GOLD)

    add_text_box(slide, 3.8, 4.95, 5.7, 0.35, 'GET IN TOUCH',
                 font_size=16, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    contact_lines = [
        'jack@thearena.protocol',
        'thearena.protocol',
        'github.com/arena-protocol',
    ]
    for j, line in enumerate(contact_lines):
        add_text_box(slide, 3.8, 5.4 + j * 0.35, 5.7, 0.3, line,
                     font_size=14, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Footer
    add_text_box(slide, 3.0, 7.0, 7.3, 0.3,
                 'Built on Base (Ethereum L2)  |  Solidity 0.8.24  |  EIP-170 Compliant',
                 font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 12)


# ═══════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════

def main():
    prs = Presentation()

    # Set widescreen 16:9
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build all 12 slides
    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_how_it_works(prs)
    slide_05_market(prs)
    slide_06_business_model(prs)
    slide_07_traction(prs)
    slide_08_financials(prs)
    slide_09_saas(prs)
    slide_10_gtm(prs)
    slide_11_use_of_funds(prs)
    slide_12_team(prs)

    output_path = '/Users/JackArnot/Desktop/The Arena Working /arena-codex/arena-pitch-deck.pptx'
    prs.save(output_path)
    print(f'Pitch deck saved to: {output_path}')
    print(f'Slides: {len(prs.slides)}')

    # Count shapes
    total_shapes = sum(len(slide.shapes) for slide in prs.slides)
    print(f'Total shapes: {total_shapes}')


if __name__ == '__main__':
    main()
